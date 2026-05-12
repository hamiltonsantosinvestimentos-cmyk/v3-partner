from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── Cores V3 ────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x09, 0x08, 0x1A)
NAVY_CARD = RGBColor(0x16, 0x27, 0x44)
NAVY_MED  = RGBColor(0x11, 0x1F, 0x35)
GOLD      = RGBColor(0xC9, 0xA8, 0x4C)
GOLD_L    = RGBColor(0xE8, 0xC9, 0x7A)
CREAM     = RGBColor(0xF0, 0xEC, 0xE4)
MUTED     = RGBColor(0x7A, 0x8F, 0xA8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
EMERALD   = RGBColor(0x34, 0xD3, 0x99)
RED       = RGBColor(0xF8, 0x71, 0x71)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # Blank layout

# ─── Helpers ─────────────────────────────────────────────────────────────────

def bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size, bold=False, color=CREAM, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb

def label(slide, text, l, t, w=3, color=GOLD):
    txt(slide, text.upper(), l, t, w, 0.3, 8, bold=True, color=color, align=PP_ALIGN.LEFT)

def divider(slide, t, color=NAVY_CARD):
    box(slide, 0.4, t, 12.53, 0.02, color)

def kpi_card(slide, l, t, w, h, valor, label_txt, cor=GOLD):
    box(slide, l, t, w, h, NAVY_CARD)
    box(slide, l, t, w, 0.06, cor)  # topo colorido
    txt(slide, valor,     l+0.15, t+0.2,  w-0.3, 0.55, 22, bold=True, color=WHITE)
    txt(slide, label_txt, l+0.15, t+0.75, w-0.3, 0.4,   9, color=MUTED)

def table_row(slide, cols, t, h=0.35, header=False):
    x = 0.4
    bg_color = NAVY_MED if header else NAVY_CARD
    widths = [3.5, 2.2, 2.2, 2.2, 2.9]
    for i, (col, w) in enumerate(zip(cols, widths)):
        box(slide, x, t, w - 0.04, h, bg_color)
        color = GOLD if header else (CREAM if i == 0 else MUTED)
        size  = 8 if header else 9
        bold  = header
        txt(slide, col, x + 0.1, t + 0.07, w - 0.2, h, size, bold=bold, color=color)
        x += w

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — CAPA
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, 13.33, 7.5, NAVY)
# Faixa ouro lateral esquerda
box(s, 0, 0, 0.12, 7.5, GOLD)
# Círculo decorativo
box(s, 8.5, -1.5, 6, 6, NAVY_CARD)

txt(s, "V3 PARTNERS",            0.4, 1.5,  7, 0.5, 11, bold=True,  color=GOLD)
txt(s, "Plano de Expansão",      0.4, 2.1,  9, 1.2, 52, bold=True,  color=CREAM)
txt(s, "MRR via Rede de Partners", 0.4, 3.4, 9, 0.6, 22, bold=False, color=MUTED)
txt(s, "2026 – 2027",            0.4, 4.1,  5, 0.5, 14, bold=False, color=MUTED)
divider(s, 5.0, GOLD)
txt(s, "Confidencial — Uso Interno · Diretoria Executiva", 0.4, 5.1, 9, 0.4, 9, color=MUTED, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — DIAGNÓSTICO
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, 0.12, 7.5, GOLD)

label(s, "Diagnóstico Atual", 0.4, 0.3)
txt(s, "Onde estamos hoje", 0.4, 0.6, 9, 0.7, 30, bold=True, color=CREAM)
divider(s, 1.4)

items = [
    ("Plataforma",            "Pronta e em produção — 10+ módulos entregues"),
    ("Partners ativos",       "< 50 parceiros (estimativa)"),
    ("MRR atual",             "< R$ 10.000/mês"),
    ("Cobrança",              "100% manual — gargalo crítico"),
    ("Funil de aquisição",    "Inexistente até o momento"),
    ("Time comercial",        "Sem SDR / Closer dedicados"),
    ("Churn",                 "Não monitorado — risco invisível"),
    ("Produto vs. mercado",   "Entregamos 10x mais valor do que cobramos"),
]

for i, (k, v) in enumerate(items):
    row = 1.55 + i * 0.62
    box(s, 0.4, row, 12.53, 0.55, NAVY_CARD)
    if i % 2 == 0:
        box(s, 0.4, row, 0.06, 0.55, GOLD)
    else:
        box(s, 0.4, row, 0.06, 0.55, NAVY_MED)
    txt(s, k, 0.6,  row + 0.12, 3.5, 0.35, 10, bold=True,  color=GOLD if i % 2 == 0 else CREAM)
    txt(s, v, 4.2,  row + 0.12, 8.5, 0.35, 10, bold=False, color=MUTED)

txt(s, "Problema central: produto pronto, receita travada por ausência de processo comercial estruturado.",
    0.4, 6.6, 12.5, 0.5, 10, bold=True, color=GOLD, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — PRECIFICAÇÃO
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, 0.12, 7.5, GOLD)

label(s, "Precificação", 0.4, 0.3)
txt(s, "Nova estrutura de planos", 0.4, 0.6, 9, 0.7, 30, bold=True, color=CREAM)
divider(s, 1.4)

planos = [
    ("STARTER",      "R$ 297/mês", "25%", "Assessores iniciantes, testando o mercado", MUTED),
    ("PARTNER",      "R$ 497/mês", "35%", "Partners ativos, 1–5 operações/mês",        GOLD),
    ("PARTNER PRO",  "R$ 897/mês", "50%", "Alta produção + co-branding exclusivo",      EMERALD),
    ("ENTERPRISE",   "R$ 2.500+",  "55%", "Escritórios, franquias, white-label total",  RGBColor(0xA7,0x8B,0xFA)),
]

for i, (nome, preco, com, desc, cor) in enumerate(planos):
    l = 0.4 + i * 3.15
    box(s, l, 1.6, 2.95, 4.2, NAVY_CARD)
    box(s, l, 1.6, 2.95, 0.07, cor)
    txt(s, nome,  l+0.15, 1.75, 2.7, 0.4, 10, bold=True,  color=cor)
    txt(s, preco, l+0.15, 2.25, 2.7, 0.7, 24, bold=True,  color=WHITE)
    txt(s, f"Comissão: {com}", l+0.15, 2.95, 2.7, 0.35, 10, bold=True, color=GOLD)
    txt(s, desc,  l+0.15, 3.4,  2.7, 0.8, 9,  bold=False, color=MUTED)

box(s, 0.4, 5.9, 12.53, 0.7, NAVY_MED)
txt(s, "Estratégia de transição:", 0.6, 6.0, 3, 0.5, 10, bold=True, color=GOLD)
txt(s, "Partners atuais mantêm preço de fundador por 6 meses. Novos partners entram no preço novo a partir de agora.",
    3.5, 6.0, 9.2, 0.5, 9, color=CREAM)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — 5 ALAVANCAS
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, 0.12, 7.5, GOLD)

label(s, "Estratégia", 0.4, 0.3)
txt(s, "As 5 Alavancas de MRR", 0.4, 0.6, 9, 0.7, 30, bold=True, color=CREAM)
divider(s, 1.4)

alavancas = [
    ("01", "AQUISIÇÃO",  "Trazer novos partners via outbound, referral e conteúdo",              GOLD),
    ("02", "ATIVAÇÃO",   "Converter trial em pagante com sequência de emails e CS ativo",         RGBColor(0x60,0xA5,0xFA)),
    ("03", "RETENÇÃO",   "Evitar churn com NPS, alertas proativos e suporte dedicado",            EMERALD),
    ("04", "EXPANSÃO",   "Upsell Starter → PRO com comparativo automático de ganhos",             RGBColor(0xA7,0x8B,0xFA)),
    ("05", "REFERRAL",   "Partner indica partner — R$150 crédito para quem indica + 1 mês grátis", RGBColor(0xF5,0x9E,0x0B)),
]

for i, (num, titulo, desc, cor) in enumerate(alavancas):
    row = 1.6 + i * 1.0
    box(s, 0.4, row, 12.53, 0.88, NAVY_CARD)
    box(s, 0.4, row, 0.06, 0.88, cor)
    txt(s, num,    0.6,  row + 0.2, 0.6, 0.5, 14, bold=True, color=cor)
    txt(s, titulo, 1.3,  row + 0.1, 2.5, 0.35, 12, bold=True, color=WHITE)
    txt(s, desc,   1.3,  row + 0.48, 11.2, 0.35, 9, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — FASE 1
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, 0.12, 7.5, GOLD)

label(s, "Fase 1 · Mês 1–2", 0.4, 0.3)
txt(s, "Fundação — Arrumar a casa antes de escalar", 0.4, 0.6, 10, 0.7, 26, bold=True, color=CREAM)
txt(s, "Meta: MRR R$ 15.000 · 35 partners ativos", 0.4, 1.15, 9, 0.4, 12, color=GOLD)
divider(s, 1.6)

acoes = [
    ("1", "Cobrança Automática",     "Integrar Asaas ou Stripe — acesso controlado por webhook de pagamento",       GOLD),
    ("2", "Time Comercial Mínimo",   "1 SDR (prospecção ativa) + 1 Closer (conversão) + 1 CS (retenção)",           RGBColor(0x60,0xA5,0xFA)),
    ("3", "Playbook Comercial",      "ICP definido, pitch de 60s, objeções mapeadas, demo roteirizada de 15 min",   EMERALD),
    ("4", "Emails de Onboarding",    "7 emails automáticos: D0 → D3 → D7 → D15 → D25 → D28 → D30",                 RGBColor(0xA7,0x8B,0xFA)),
    ("5", "Ajuste de Precificação",  "Novos planos: Starter R$297 · Partner R$497 · PRO R$897 · Enterprise R$2.5k", RGBColor(0xF5,0x9E,0x0B)),
]

for i, (n, titulo, desc, cor) in enumerate(acoes):
    row = 1.75 + i * 0.98
    box(s, 0.4, row, 12.53, 0.86, NAVY_CARD)
    txt(s, n,      0.6,  row+0.22, 0.4,  0.4, 16, bold=True, color=cor)
    txt(s, titulo, 1.1,  row+0.1,  3.5,  0.35, 11, bold=True, color=WHITE)
    txt(s, desc,   1.1,  row+0.48, 11.5, 0.32,  9, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — CANAIS DE AQUISIÇÃO
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, 0.12, 7.5, GOLD)

label(s, "Fase 2 · Mês 3–6", 0.4, 0.3)
txt(s, "Canais de Aquisição", 0.4, 0.6, 9, 0.7, 30, bold=True, color=CREAM)
txt(s, "Meta: 150 partners · MRR R$ 75.000", 0.4, 1.15, 9, 0.4, 12, color=GOLD)
divider(s, 1.6)

canais = [
    ("LinkedIn Outbound",     "SDR prospecta 30 leads/dia — assessores XP/BTG, contadores, consultores financeiros",   "5 demos/sem → 2 trials → 0,6 fechamentos", GOLD),
    ("Programa de Referral",  "Link exclusivo por partner — R$150 crédito + 1 mês grátis para indicado",               "Meta: 1 indicação a cada 3 partners/mês",  RGBColor(0x60,0xA5,0xFA)),
    ("Parceiros Estratégicos","Escritórios contábeis, corretoras, sindicatos (ANCORD, ABAI, CRC)",                      "Palestra gratuita → converte 5-10% trial",  EMERALD),
    ("Conteúdo + SEO",        "YouTube e LinkedIn pessoal dos sócios — 1 post/dia, 1 vídeo/semana",                    "Meta: 2-3 trials/semana via inbound",       RGBColor(0xA7,0x8B,0xFA)),
]

for i, (canal, desc, meta, cor) in enumerate(canais):
    row = 1.75 + i * 1.32
    box(s, 0.4,  row, 8.6, 1.18, NAVY_CARD)
    box(s, 9.15, row, 4.0, 1.18, NAVY_MED)
    box(s, 0.4,  row, 0.06, 1.18, cor)
    txt(s, canal, 0.6,  row+0.1,  8.2, 0.35, 11, bold=True,  color=cor)
    txt(s, desc,  0.6,  row+0.52, 8.2, 0.55,  9, color=MUTED)
    txt(s, "META", 9.25, row+0.1,  3.5, 0.25,  7, bold=True, color=MUTED)
    txt(s, meta,  9.25, row+0.38, 3.7, 0.65,  9, bold=True,  color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — PROJEÇÃO FINANCEIRA
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, 0.12, 7.5, GOLD)

label(s, "Projeção Financeira", 0.4, 0.3)
txt(s, "MRR · 12 meses", 0.4, 0.6, 9, 0.7, 30, bold=True, color=CREAM)
divider(s, 1.4)

# KPIs topo
kpi_card(s, 0.4,  1.55, 2.9, 1.2, "R$ 200k",   "MRR no Mês 12",         GOLD)
kpi_card(s, 3.45, 1.55, 2.9, 1.2, "400",        "Partners Ativos (M12)",  RGBColor(0x60,0xA5,0xFA))
kpi_card(s, 6.5,  1.55, 2.9, 1.2, "R$ 3,1 MM", "ARR Projetado",          EMERALD)
kpi_card(s, 9.55, 1.55, 2.9, 1.2, "< 3%",      "Churn Mensal Alvo",      RGBColor(0xA7,0x8B,0xFA))

# Tabela
headers = ["MÊS / META", "PARTNERS", "MRR ASSIN.", "MRR TOTAL", "ACUMULADO"]
table_row(s, headers, 3.0, h=0.38, header=True)

rows_data = [
    ("Mês 2",  "35",    "R$ 14k",   "R$ 15k",   "R$ 23k"),
    ("Mês 3",  "55",    "R$ 22k",   "R$ 24k",   "R$ 47k"),
    ("Mês 6",  "150",   "R$ 60k",   "R$ 75k",   "R$ 208k"),
    ("Mês 9",  "250",   "R$ 110k",  "R$ 140k",  "R$ 560k"),
    ("Mês 12", "400",   "R$ 200k",  "R$ 260k",  "R$ 1,2 MM"),
]

for i, row in enumerate(rows_data):
    t = 3.4 + i * 0.62
    table_row(s, row, t, h=0.56)

box(s, 0.4, 6.7, 12.53, 0.55, NAVY_MED)
txt(s, "Break-even comercial:", 0.6, 6.8, 2.8, 0.4, 10, bold=True, color=GOLD)
txt(s, "~60 partners (MRR cobre custo total do time comercial)", 3.5, 6.8, 9, 0.4, 10, color=CREAM)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — FASE 3 / ESCALA
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, 0.12, 7.5, GOLD)

label(s, "Fase 3 · Mês 7–12", 0.4, 0.3)
txt(s, "Escala — Multiplicar o que funciona", 0.4, 0.6, 9, 0.7, 30, bold=True, color=CREAM)
txt(s, "Meta: 400 partners · MRR R$ 200.000", 0.4, 1.15, 9, 0.4, 12, color=GOLD)
divider(s, 1.6)

itens = [
    ("Expansão Geográfica",   "SP (base) → MG, PR, GO (mês 4) → Nordeste e Sul (mês 8)\nPartner Âncora Regional: 1 PRO por cidade com comissão extra por indicações locais", RGBColor(0x60,0xA5,0xFA)),
    ("Tier Enterprise",       "R$2.500–5.000/mês · 55% comissão · Multi-usuário · White-label · SLA dedicado\nPerfil: escritórios com 5+ assessores, family offices, corretoras", GOLD),
    ("Academy Monetizado",    "Certificação 'V3 Certified Partner' R$497 (único) + trilhas avançadas R$97/mês\n20% dos partners = +R$20.000 MRR adicional", EMERALD),
    ("Marketplace Recorrente","Listing fee: R$150/mês por fornecedor · Destaque R$300/mês\nMeta: 30 fornecedores = R$4.500–9.000/mês recorrente", RGBColor(0xA7,0x8B,0xFA)),
]

for i, (titulo, desc, cor) in enumerate(itens):
    row = 1.75 + i * 1.3
    box(s, 0.4, row, 12.53, 1.15, NAVY_CARD)
    box(s, 0.4, row, 0.06, 1.15, cor)
    txt(s, titulo, 0.65, row+0.1,  5,    0.35, 11, bold=True, color=cor)
    txt(s, desc,   0.65, row+0.5,  12.1, 0.6,   9, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — ROADMAP PLATAFORMA
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, 0.12, 7.5, GOLD)

label(s, "Plataforma", 0.4, 0.3)
txt(s, "Roadmap de Implementação Técnica", 0.4, 0.6, 10, 0.7, 30, bold=True, color=CREAM)
divider(s, 1.4)

roadmap = [
    ("1", "Cobrança Automática (Asaas/Stripe)",        "CRÍTICO",     "Semana 1–2",    GOLD),
    ("2", "Emails automáticos de onboarding (7 emails)", "ALTO",      "Semana 2–3",    RGBColor(0x60,0xA5,0xFA)),
    ("3", "Programa de Referral com link rastreável",   "ALTO",        "Mês 2",         EMERALD),
    ("4", "Alertas de churn (trial vencendo)",          "ALTO",        "Mês 2",         EMERALD),
    ("5", "Upgrade incentivado no dashboard",           "MÉDIO",       "Mês 3",         RGBColor(0xA7,0x8B,0xFA)),
    ("6", "Multi-usuário Enterprise",                   "MÉDIO",       "Mês 5",         RGBColor(0xA7,0x8B,0xFA)),
    ("7", "Academy monetizado + certificação",          "MÉDIO",       "Mês 8",         MUTED),
]

for i, (n, feat, impacto, prazo, cor) in enumerate(roadmap):
    row = 1.55 + i * 0.72
    box(s, 0.4, row, 12.53, 0.64, NAVY_CARD)
    box(s, 0.4, row, 0.06, 0.64, cor)
    txt(s, n,      0.65, row+0.17, 0.4,  0.35,  12, bold=True, color=cor)
    txt(s, feat,   1.15, row+0.17, 7.5,  0.35,  10, bold=True, color=WHITE)
    # Badge impacto
    bcor = GOLD if impacto == "CRÍTICO" else (EMERALD if impacto == "ALTO" else MUTED)
    txt(s, impacto, 9.0, row+0.17, 1.8, 0.35, 8, bold=True, color=bcor)
    txt(s, prazo,   10.9, row+0.17, 1.8, 0.35, 9, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — GOVERNANÇA
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, 0.12, 7.5, GOLD)

label(s, "Governança", 0.4, 0.3)
txt(s, "Cadência de Gestão e KPIs", 0.4, 0.6, 9, 0.7, 30, bold=True, color=CREAM)
divider(s, 1.4)

# Reuniões
txt(s, "CADÊNCIA DE REUNIÕES", 0.4, 1.55, 6, 0.3, 8, bold=True, color=GOLD)

reunioes = [
    ("War Room Comercial",   "Semanal (segunda)",   "Pipeline, demos, fechamentos, blockers"),
    ("Review de MRR",        "Quinzenal",            "Dashboard métricas, churn, campanhas"),
    ("Reunião de Diretoria", "Mensal",               "Resultado vs meta, decisões estratégicas"),
    ("Partner Advisory",     "Trimestral",           "Feedback de partners, NPS, melhorias"),
]

for i, (nome, freq, pauta) in enumerate(reunioes):
    row = 1.95 + i * 0.85
    box(s, 0.4, row, 7.0, 0.75, NAVY_CARD)
    txt(s, nome,  0.6, row+0.08, 3.0, 0.3, 10, bold=True, color=CREAM)
    txt(s, freq,  0.6, row+0.42, 2.0, 0.28, 8, bold=True, color=GOLD)
    txt(s, pauta, 3.5, row+0.2,  3.7, 0.4,  9, color=MUTED)

# KPIs
txt(s, "KPIS SEMANAIS", 7.7, 1.55, 5, 0.3, 8, bold=True, color=GOLD)

kpis = [
    ("MRR",                "Crescer 20%/mês"),
    ("Churn Rate",         "< 3%/mês"),
    ("Trial → Pago",       "> 30% conversão"),
    ("NPS",                "> 50"),
    ("Novos Leads/sem",    "> 30"),
    ("Demos realizadas",   "> 4/semana"),
]

for i, (k, v) in enumerate(kpis):
    row = 1.95 + i * 0.72
    box(s, 7.7, row, 5.03, 0.64, NAVY_CARD)
    txt(s, k, 7.9, row+0.1,  2.5, 0.3, 9, bold=True, color=CREAM)
    txt(s, v, 7.9, row+0.38, 4.6, 0.22, 8, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — PRÓXIMOS PASSOS
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, 0.12, 7.5, GOLD)

label(s, "Decisão", 0.4, 0.3)
txt(s, "3 Ações Imediatas da Diretoria", 0.4, 0.6, 10, 0.7, 30, bold=True, color=CREAM)
divider(s, 1.4)

acoes_dir = [
    ("01", "Integrar cobrança automática",
     "Asaas ou Stripe — sem isso nenhuma escala é sustentável.\nResponsável: Hamilton · Prazo: 2 semanas",
     GOLD),
    ("02", "Definir o time comercial",
     "SDR + Closer — pode ser um dos sócios por 60 dias enquanto contratação ocorre.\nResponsável: João · Prazo: esta semana",
     RGBColor(0x60,0xA5,0xFA)),
    ("03", "Ajustar a precificação",
     "Migrar para novos planos e comunicar 'Preço de Fundador' para base atual.\nResponsável: Robson + Hamilton · Prazo: 1 semana",
     EMERALD),
]

for i, (n, titulo, desc, cor) in enumerate(acoes_dir):
    row = 1.6 + i * 1.55
    box(s, 0.4, row, 12.53, 1.38, NAVY_CARD)
    box(s, 0.4, row, 0.1, 1.38, cor)
    txt(s, n,      0.7,  row+0.3,  0.8, 0.6, 28, bold=True, color=cor)
    txt(s, titulo, 1.7,  row+0.15, 10,  0.45, 16, bold=True, color=WHITE)
    txt(s, desc,   1.7,  row+0.65, 10.8, 0.65, 10, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — ENCERRAMENTO
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, 0.12, 7.5, GOLD)
box(s, 0, 6.0, 13.33, 1.5, NAVY_MED)

txt(s, "V3 PARTNERS",                             0.4, 1.8, 12, 0.6, 14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(s, "O produto está pronto.",                   0.4, 2.5, 12, 0.8, 44, bold=True, color=CREAM, align=PP_ALIGN.CENTER)
txt(s, "Agora é hora de escalar.",                 0.4, 3.3, 12, 0.8, 44, bold=True, color=GOLD,  align=PP_ALIGN.CENTER)
divider(s, 4.3)
txt(s, "ARR alvo em 12 meses: R$ 3.100.000",       0.4, 4.5, 12, 0.5, 16, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, "v3partners.com.br  ·  Confidencial",       0.4, 6.1, 12, 0.4, 10, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SALVAR
# ─────────────────────────────────────────────────────────────────────────────
output = r"C:\Users\atend\Desktop\V3Partners_Plano_Expansao_MRR.pptx"
prs.save(output)
print(f"✓ Salvo em: {output}")
